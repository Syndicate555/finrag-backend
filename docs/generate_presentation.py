from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(0x0D, 0x1B, 0x2A)
DARK_NAVY = RGBColor(0x09, 0x12, 0x1C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_BLUE = RGBColor(0x00, 0x9E, 0xDB)
ACCENT_ORANGE = RGBColor(0xEB, 0x8C, 0x00)
MUTED_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
SOFT_WHITE = RGBColor(0xE8, 0xEC, 0xF1)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
PURPLE = RGBColor(0x7C, 0x4D, 0xFF)
TEAL = RGBColor(0x00, 0x96, 0x88)
RED_ACCENT = RGBColor(0xE5, 0x39, 0x35)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title_text):
    # Title bar background
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_NAVY
    bar.line.fill.background()

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.1), prs.slide_width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.85))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"

    # "FinRAG" branding top-right
    brand = slide.shapes.add_textbox(Inches(11), Inches(0.25), Inches(2), Inches(0.6))
    bf = brand.text_frame
    bp = bf.paragraphs[0]
    bp.alignment = PP_ALIGN.RIGHT
    run = bp.add_run()
    run.text = "FinRAG"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = ACCENT_BLUE
    run.font.name = "Calibri"


def add_box(slide, left, top, width, height, fill_color, border_color, texts, font_size=14):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = border_color
    box.line.width = Pt(1.5)
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    for i, (text, bold, color, size_override) in enumerate(texts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size_override or font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"

    return box


def add_arrow(slide, start_left, start_top, end_left, end_top):
    connector = slide.shapes.add_connector(
        1, start_left, start_top, end_left, end_top  # straight connector
    )
    connector.line.color.rgb = ACCENT_BLUE
    connector.line.width = Pt(2)


def add_bullet_text(slide, left, top, width, height, items, font_size=16, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, indent_level, bold, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = indent_level
        p.space_after = spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"

    return txBox


# ─────────────────────────────────────────────
# SLIDE 1: Problem Statement & Solution Overview
# ─────────────────────────────────────────────
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1, NAVY)
add_title_bar(slide1, "Problem Statement & Solution Overview")

# Problem section
add_bullet_text(slide1, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.8), [
    ("THE PROBLEM", 0, True, ACCENT_ORANGE),
    ("Financial analysts need to extract insights from complex documents containing text, tables, and images", 0, False, WHITE),
    ("Manual review is time-consuming and error-prone", 0, False, LIGHT_GRAY),
    ("Existing tools lack contextual understanding and citation support", 0, False, LIGHT_GRAY),
], font_size=16, spacing=Pt(10))

# Solution section
add_bullet_text(slide1, Inches(0.6), Inches(4.0), Inches(5.8), Inches(3.2), [
    ("THE SOLUTION: FinRAG", 0, True, ACCENT_BLUE),
    ("RAG-powered conversational AI for financial document analysis", 0, False, WHITE),
    ("Intelligent query routing: KB queries, general conversation, and clarification requests", 0, False, LIGHT_GRAY),
    ("Citations with page numbers and relevance scores for every response", 0, False, LIGHT_GRAY),
    ("Real-time streaming responses via Server-Sent Events", 0, False, LIGHT_GRAY),
], font_size=16, spacing=Pt(10))

# Key capabilities boxes on the right
caps = [
    ("Knowledge Base\nConversations", ACCENT_BLUE),
    ("General\nConversations", TEAL),
    ("Citations with\nPage & Relevance", ACCENT_ORANGE),
    ("Real-time\nStreaming (SSE)", PURPLE),
]
for i, (cap_text, color) in enumerate(caps):
    y = Inches(1.5) + Inches(1.4) * i
    add_box(slide1, Inches(7.2), y, Inches(2.8), Inches(1.1), MUTED_BLUE, color, [
        (cap_text, True, WHITE, 15),
    ])

# Tech badges on far right
techs = ["GPT-4o", "Pinecone", "FastAPI", "Next.js"]
for i, tech in enumerate(techs):
    y = Inches(1.55) + Inches(1.4) * i
    badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.3), y, Inches(1.6), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_NAVY
    badge.line.color.rgb = ACCENT_BLUE
    badge.line.width = Pt(1)
    tf = badge.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = tech
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = ACCENT_BLUE
    run.font.name = "Calibri"

# ─────────────────────────────────────────────
# SLIDE 2: Solution Architecture
# ─────────────────────────────────────────────
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, NAVY)
add_title_bar(slide2, "Solution Architecture — Cloud Resources")

# Architecture boxes - row 1 (user-facing)
components_row1 = [
    ("Frontend", "Next.js 16 on Vercel\nfinrag.info", ACCENT_BLUE),
    ("Backend", "FastAPI on Azure\nContainer Apps\n(serverless, scale-to-zero)", TEAL),
    ("LLM", "OpenAI GPT-4o\n(generation)\nGPT-4o-mini\n(query classification)", ACCENT_ORANGE),
]
x_positions_r1 = [Inches(0.4), Inches(4.2), Inches(8.8)]
for (title, desc, color), x in zip(components_row1, x_positions_r1):
    add_box(slide2, x, Inches(1.5), Inches(3.5), Inches(1.8), MUTED_BLUE, color, [
        (title, True, color, 16),
        (desc, False, LIGHT_GRAY, 12),
    ])

# Row 2 (data layer)
components_row2 = [
    ("Vector DB", "Pinecone\ntext-embedding-3-large\n1536 dimensions", PURPLE),
    ("Database", "Supabase PostgreSQL\nthreads, messages,\ndocuments, sections, feedback", GREEN),
    ("Storage", "Supabase Object Storage\nPDF blob storage", RGBColor(0xFF, 0xB7, 0x4D)),
]
x_positions_r2 = [Inches(0.4), Inches(4.2), Inches(8.8)]
for (title, desc, color), x in zip(components_row2, x_positions_r2):
    add_box(slide2, x, Inches(3.7), Inches(3.5), Inches(1.6), MUTED_BLUE, color, [
        (title, True, color, 16),
        (desc, False, LIGHT_GRAY, 12),
    ])

# Row 3 (infrastructure)
components_row3 = [
    ("Document Intelligence", "Azure Document Intelligence\nStructured parsing\n(fallback: pdfplumber)", RED_ACCENT),
    ("CI/CD", "GitHub Actions\n→ ACR → Container Apps", RGBColor(0x90, 0xA4, 0xAE)),
    ("Monitoring", "Azure Application Insights\n+ Azure Portal Dashboard", RGBColor(0x42, 0xA5, 0xF5)),
]
x_positions_r3 = [Inches(0.4), Inches(4.2), Inches(8.8)]
for (title, desc, color), x in zip(components_row3, x_positions_r3):
    add_box(slide2, x, Inches(5.7), Inches(3.5), Inches(1.4), MUTED_BLUE, color, [
        (title, True, color, 16),
        (desc, False, LIGHT_GRAY, 12),
    ])

# Arrows between row 1 boxes
add_arrow(slide2, Inches(3.9), Inches(2.4), Inches(4.2), Inches(2.4))
add_arrow(slide2, Inches(7.7), Inches(2.4), Inches(8.8), Inches(2.4))

# Arrows from row 1 to row 2
add_arrow(slide2, Inches(5.95), Inches(3.3), Inches(5.95), Inches(3.7))
add_arrow(slide2, Inches(2.15), Inches(3.3), Inches(2.15), Inches(3.7))
add_arrow(slide2, Inches(10.55), Inches(3.3), Inches(10.55), Inches(3.7))

# ─────────────────────────────────────────────
# SLIDE 3: Design Flow (RAG Pipeline)
# ─────────────────────────────────────────────
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, NAVY)
add_title_bar(slide3, "Design Flow — RAG Pipeline")

# Ingestion Pipeline section
add_bullet_text(slide3, Inches(0.4), Inches(1.3), Inches(6), Inches(0.5), [
    ("INGESTION PIPELINE", 0, True, ACCENT_ORANGE),
], font_size=18)

ingestion_steps = [
    ("Upload PDF", ACCENT_BLUE),
    ("Azure DI Parse\n(fallback:\npdfplumber)", RED_ACCENT),
    ("Token Chunking\n512 tokens\n64 overlap", TEAL),
    ("Embed\ntext-embedding\n-3-large", PURPLE),
    ("Upsert to\nPinecone +\nmetadata", GREEN),
]
for i, (text, color) in enumerate(ingestion_steps):
    x = Inches(0.4) + Inches(2.5) * i
    add_box(slide3, x, Inches(1.9), Inches(2.2), Inches(1.3), MUTED_BLUE, color, [
        (text, True, WHITE, 13),
    ])
    if i < len(ingestion_steps) - 1:
        add_arrow(slide3, x + Inches(2.2), Inches(2.55), x + Inches(2.5), Inches(2.55))

# Query Pipeline section
add_bullet_text(slide3, Inches(0.4), Inches(3.5), Inches(6), Inches(0.5), [
    ("QUERY PIPELINE", 0, True, ACCENT_BLUE),
], font_size=18)

query_steps = [
    ("User\nMessage", ACCENT_BLUE),
    ("Query Router\nkb / general /\nclarification", ACCENT_ORANGE),
    ("Retrieve\ntop-k=20 from\nPinecone", PURPLE),
    ("Stream GPT-4o\nvia SSE with\ncitations", TEAL),
    ("Display with\npage refs &\nrelevance", GREEN),
]
for i, (text, color) in enumerate(query_steps):
    x = Inches(0.4) + Inches(2.5) * i
    add_box(slide3, x, Inches(4.1), Inches(2.2), Inches(1.3), MUTED_BLUE, color, [
        (text, True, WHITE, 13),
    ])
    if i < len(query_steps) - 1:
        add_arrow(slide3, x + Inches(2.2), Inches(4.75), x + Inches(2.5), Inches(4.75))

# Metadata tags section
add_bullet_text(slide3, Inches(0.4), Inches(5.7), Inches(12.5), Inches(1.5), [
    ("METADATA TAGS", 0, True, RGBColor(0x90, 0xA4, 0xAE)),
    ("section_heading  |  section_level  |  parent_section  |  content_type (text/table)  |  page_start  |  page_end", 0, False, LIGHT_GRAY),
], font_size=14, spacing=Pt(6))


# ─────────────────────────────────────────────
# SLIDE 4: Key Decisions & Evaluation
# ─────────────────────────────────────────────
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, NAVY)
add_title_bar(slide4, "Key Decisions & Evaluation Approach")

# Left column: Key Decisions
add_bullet_text(slide4, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5), [
    ("KEY ARCHITECTURAL DECISIONS", 0, True, ACCENT_ORANGE),
], font_size=18)

decisions = [
    ("Azure Container Apps over VM — serverless, scale-to-zero, cost-efficient", 0, False, WHITE),
    ("Pinecone over pgvector — managed service, better scale, metadata filtering", 0, False, LIGHT_GRAY),
    ("Token-based chunking (512 tokens, 64 overlap) — preserves semantic boundaries", 0, False, WHITE),
    ("Azure DI with pdfplumber fallback — robust parsing for tables + text", 0, False, LIGHT_GRAY),
    ("GPT-4o-mini for query routing — fast classification, cost-effective", 0, False, WHITE),
    ("SSE streaming — real-time UX, lower perceived latency", 0, False, LIGHT_GRAY),
    ("Supabase for DB + Storage — unified platform, real-time capabilities", 0, False, WHITE),
]
add_bullet_text(slide4, Inches(0.5), Inches(1.9), Inches(6.2), Inches(5.2), decisions, font_size=14, spacing=Pt(10))

# Right column: Evaluation
add_bullet_text(slide4, Inches(7.0), Inches(1.3), Inches(6), Inches(0.5), [
    ("EVALUATION APPROACH", 0, True, ACCENT_BLUE),
], font_size=18)

# Offline eval box
add_box(slide4, Inches(7.0), Inches(1.9), Inches(5.8), Inches(2.4), MUTED_BLUE, PURPLE, [
    ("Offline — RAGAS Framework", True, PURPLE, 16),
    ("", False, WHITE, 6),
    ("Faithfulness — Is the answer grounded in context?", False, LIGHT_GRAY, 13),
    ("Answer Relevancy — Does it address the query?", False, LIGHT_GRAY, 13),
    ("Context Precision — Are retrieved chunks relevant?", False, LIGHT_GRAY, 13),
    ("Context Recall — Are all needed chunks retrieved?", False, LIGHT_GRAY, 13),
])

# Online eval box
add_box(slide4, Inches(7.0), Inches(4.6), Inches(5.8), Inches(2.4), MUTED_BLUE, GREEN, [
    ("Online — Production Monitoring", True, GREEN, 16),
    ("", False, WHITE, 6),
    ("User Feedback — Thumbs up/down per response", False, LIGHT_GRAY, 13),
    ("Application Insights — Latency, error rate tracking", False, LIGHT_GRAY, 13),
    ("Azure Dashboard — Dependency duration, throughput", False, LIGHT_GRAY, 13),
    ("Usage Analytics — Query patterns, document coverage", False, LIGHT_GRAY, 13),
])


# ─────────────────────────────────────────────
# SLIDE 5: Live Demo & Future Enhancements
# ─────────────────────────────────────────────
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, NAVY)
add_title_bar(slide5, "Live Demo & Future Enhancements")

# Left: Demo walkthrough
add_bullet_text(slide5, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5), [
    ("DEMO WALKTHROUGH", 0, True, ACCENT_ORANGE),
], font_size=18)

demo_steps = [
    ("1.", "Upload BMO Annual Report (PDF) and observe ingestion pipeline", ACCENT_BLUE),
    ("2.", "Ask a KB question — see citations with page numbers & relevance", TEAL),
    ("3.", "Ask a general question — observe query routing classification", PURPLE),
    ("4.", "Send an ambiguous query — trigger clarification flow", ACCENT_ORANGE),
    ("5.", "Submit feedback — thumbs up/down stored in Supabase", GREEN),
]
for i, (num, text, color) in enumerate(demo_steps):
    y = Inches(2.0) + Inches(0.9) * i
    # Number circle
    circle = slide5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    # Step text
    txBox = slide5.shapes.add_textbox(Inches(1.3), y, Inches(5.2), Inches(0.5))
    stf = txBox.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.LEFT
    srun = sp.add_run()
    srun.text = text
    srun.font.size = Pt(15)
    srun.font.color.rgb = WHITE
    srun.font.name = "Calibri"

# Right: Future Enhancements
add_bullet_text(slide5, Inches(7.0), Inches(1.3), Inches(6), Inches(0.5), [
    ("FUTURE ENHANCEMENTS", 0, True, ACCENT_BLUE),
], font_size=18)

future_items = [
    ("Multi-document cross-referencing and comparison", GREEN),
    ("Advanced table extraction with structure preservation", TEAL),
    ("Image/chart analysis with vision models", PURPLE),
    ("Hybrid search (dense + sparse retrieval)", ACCENT_BLUE),
    ("Fine-tuned embedding model for financial domain", ACCENT_ORANGE),
    ("Role-based access control (RBAC)", RED_ACCENT),
    ("Automated RAGAS evaluation pipeline in CI/CD", RGBColor(0x90, 0xA4, 0xAE)),
]
for i, (text, color) in enumerate(future_items):
    y = Inches(2.0) + Inches(0.7) * i
    # Accent bar
    bar = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), y + Pt(2), Pt(4), Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    txBox = slide5.shapes.add_textbox(Inches(7.3), y, Inches(5.5), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(15)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

# ─── Footer on all slides ───
for slide in prs.slides:
    footer = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(4), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "PwC Take-Home Assessment — Saffat Aziz"
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x66, 0x77, 0x88)
    run.font.name = "Calibri"

prs.save("docs/FinRAG_Presentation.pptx")
print("Presentation saved to docs/FinRAG_Presentation.pptx")
print(f"Total slides: {len(prs.slides)}")
